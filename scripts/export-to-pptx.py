#!/usr/bin/env python3
"""
범용 마크다운/텍스트 → PPTX 변환 스크��트

사용법:
    python scripts/export-to-pptx.py <input.md> [output.pptx] [--theme dark|light|brand]

input:  마크다운 파일 경로
output: PPTX 출력 경로 (생략 시 input과 같은 위치에 .pptx로 저장)

예시:
    python scripts/export-to-pptx.py clients/huenic/proposals/2026-04-13-proposal.md
    python scripts/export-to-pptx.py report.md report.pptx --theme dark
"""
import sys
import re
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ── Theme Definitions ──
THEMES = {
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x1A, 0x1A),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "sub_color": RGBColor(0x66, 0x66, 0x66),
        "accent": RGBColor(0x00, 0x66, 0xFF),
        "divider": RGBColor(0xE0, 0xE0, 0xE0),
    },
    "dark": {
        "bg": RGBColor(0x1A, 0x1B, 0x18),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0xCC, 0xCC, 0xCC),
        "sub_color": RGBColor(0x99, 0x99, 0x99),
        "accent": RGBColor(0xA2, 0xC7, 0xE2),
        "divider": RGBColor(0x33, 0x33, 0x33),
    },
    "brand": {
        "bg": RGBColor(0xFA, 0xF6, 0xF1),
        "title_color": RGBColor(0x1A, 0x1A, 0x1A),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "sub_color": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0xE8, 0x65, 0x2B),
        "divider": RGBColor(0xE8, 0xE0, 0xD6),
    },
}

# ── Slide Dimensions (16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_T = Inches(0.6)
CONTENT_W = Inches(11.7)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def parse_markdown(md_text):
    """Parse markdown into slide-friendly sections."""
    slides = []
    current_slide = None
    current_bullets = []

    for line in md_text.split('\n'):
        line = line.rstrip()

        # H1 = Title slide
        if line.startswith('# ') and not line.startswith('## '):
            if current_slide:
                current_slide['bullets'] = current_bullets
                slides.append(current_slide)
            current_slide = {'type': 'title', 'title': line[2:].strip(), 'bullets': []}
            current_bullets = []

        # H2 = Section slide
        elif line.startswith('## '):
            if current_slide:
                current_slide['bullets'] = current_bullets
                slides.append(current_slide)
            current_slide = {'type': 'section', 'title': line[3:].strip(), 'bullets': []}
            current_bullets = []

        # H3 = Sub-section (add as bold bullet)
        elif line.startswith('### '):
            if not current_slide:
                current_slide = {'type': 'section', 'title': '', 'bullets': []}
            current_bullets.append({'text': line[4:].strip(), 'bold': True, 'level': 0})

        # Bullets
        elif line.startswith('- ') or line.startswith('* '):
            if not current_slide:
                current_slide = {'type': 'section', 'title': '', 'bullets': []}
            text = line[2:].strip()
            # Clean markdown bold
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            current_bullets.append({'text': text, 'bold': False, 'level': 0})

        elif line.startswith('  - ') or line.startswith('  * '):
            if not current_slide:
                current_slide = {'type': 'section', 'title': '', 'bullets': []}
            text = line[4:].strip()
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            current_bullets.append({'text': text, 'bold': False, 'level': 1})

    if current_slide:
        current_slide['bullets'] = current_bullets
        slides.append(current_slide)

    return slides


def create_pptx(slides_data, theme_name="light"):
    theme = THEMES.get(theme_name, THEMES["light"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for s in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, theme["bg"])

        if s['type'] == 'title':
            # Title slide - centered
            add_text(slide, 1.5, 2.5, 10.3, 1.5, s['title'],
                     size=36, bold=True, color=theme["title_color"], align=PP_ALIGN.CENTER)
            if s['bullets']:
                subtitle = '\n'.join(b['text'] for b in s['bullets'][:3])
                add_text(slide, 2, 4.2, 9.3, 2, subtitle,
                         size=18, color=theme["sub_color"], align=PP_ALIGN.CENTER)
        else:
            # Section slide
            if s['title']:
                add_text(slide, 0.8, 0.5, 11.7, 0.8, s['title'],
                         size=28, bold=True, color=theme["title_color"])
                # Accent line
                line = slide.shapes.add_shape(
                    1, MARGIN_L, Inches(1.2), Inches(2), Pt(3))
                line.fill.solid()
                line.fill.fore_color.rgb = theme["accent"]
                line.line.fill.background()
                y_start = 1.5
            else:
                y_start = 0.6

            # Bullets
            if s['bullets']:
                y = y_start
                for bullet in s['bullets']:
                    indent = 1.2 if bullet['level'] > 0 else 0.8
                    prefix = "  " if bullet['level'] > 0 else ""
                    size = 14 if bullet['level'] > 0 else 16
                    color = theme["sub_color"] if bullet['level'] > 0 else theme["text_color"]

                    text = f"{prefix}{bullet['text']}"
                    add_text(slide, indent, y, 11.0 - indent + 0.8, 0.5, text,
                             size=size, bold=bullet.get('bold', False), color=color)
                    y += 0.45 if bullet['level'] > 0 else 0.5

                    # Page break if too many bullets
                    if y > 6.5:
                        break

    return prs


def main():
    parser = argparse.ArgumentParser(description="Markdown to PPTX converter")
    parser.add_argument("input", help="Input markdown file")
    parser.add_argument("output", nargs="?", help="Output PPTX file")
    parser.add_argument("--theme", choices=["light", "dark", "brand"], default="light")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"File not found: {input_path}")
        sys.exit(1)

    output_path = Path(args.output) if args.output else input_path.with_suffix('.pptx')

    md_text = input_path.read_text(encoding='utf-8')
    slides = parse_markdown(md_text)

    if not slides:
        print("No slides found. Make sure the markdown has # or ## headings.")
        sys.exit(1)

    prs = create_pptx(slides, args.theme)
    prs.save(str(output_path))
    print(f"PPTX exported: {output_path} ({len(slides)} slides, theme: {args.theme})")


if __name__ == "__main__":
    main()
