"""미례국밥 제안서 PPTX 생성 + Google Drive 업로드 스크립트"""

import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(__file__))

# ─── 색상 팔레트 ───
BRAND_BLUE = RGBColor(0x1A, 0x56, 0xDB)
BRAND_NAVY = RGBColor(0x10, 0x1A, 0x33)
BRAND_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF4)
ACCENT_BLUE_BG = RGBColor(0xED, 0xF2, 0xFF)
WARM_BG = RGBColor(0xFF, 0xF3, 0xED)
CARD_BG = RGBColor(0xF8, 0xF8, 0xFA)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "meta-ad-pipeline", "clients", "mirye-gukbap", "proposals")
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "20260308_mirye_proposal.pptx")

# ─── 헬퍼 함수 ───

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def set_text(shape, text, font_size=11, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스에 단일 스타일 텍스트 설정"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(text_frame, text, font_size=11, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕", space_before=0):
    """텍스트 프레임에 새 문단 추가"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """배경 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """둥근 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def multiline_textbox(slide, left, top, width, height, lines, default_size=11, default_color=BLACK):
    """여러 줄 텍스트 (각 줄별 스타일 가능)"""
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line_spec in enumerate(lines):
        if isinstance(line_spec, str):
            text, size, bold, color = line_spec, default_size, False, default_color
        else:
            text = line_spec.get("text", "")
            size = line_spec.get("size", default_size)
            bold = line_spec.get("bold", False)
            color = line_spec.get("color", default_color)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.alignment = line_spec.get("align", PP_ALIGN.LEFT) if isinstance(line_spec, dict) else PP_ALIGN.LEFT

    return tb


# ─── 슬라이드 빌드 ───

def build_presentation():
    prs = Presentation()
    # 16:9 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 빈 레이아웃 사용
    blank_layout = prs.slide_layouts[6]

    # ═══════════════════════════════════════════
    # 슬라이드 1: 표지
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, BRAND_NAVY)

    # 상단 액센트 라인
    add_rect(slide, 2, 1.8, 9.333, 0.05, BRAND_ORANGE)

    # 타이틀
    tb = add_textbox(slide, 2, 2.1, 9.333, 1.5)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "미례국밥"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "마케팅 제안서"
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    tb = add_textbox(slide, 2, 4.0, 9.333, 0.6)
    set_text(tb, "\"돼지우동, 서울이 알게 합니다\"", 22, color=BRAND_ORANGE, alignment=PP_ALIGN.CENTER)

    # 하단 라인
    add_rect(slide, 2, 5.2, 9.333, 0.05, BRAND_ORANGE)

    # 하단 정보
    tb = add_textbox(slide, 2, 5.5, 9.333, 0.8)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "2026.03  |  브랜드라이즈"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "brandrise.co.kr"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 2: 현재 상황 진단
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # 라벨
    tb = add_textbox(slide, 0.8, 0.4, 3, 0.35)
    set_text(tb, "01  현재 상황", 12, bold=True, color=BRAND_BLUE)

    # 타이틀
    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "부산에서는 성공했습니다. 문제는 서울입니다.", 30, bold=True, color=BLACK)

    # 좌측 - 잘 하고 있는 것
    add_rounded_rect(slide, 0.8, 2.0, 5.5, 2.8, ACCENT_BLUE_BG)
    tb = add_textbox(slide, 1.1, 2.1, 5, 0.4)
    set_text(tb, "잘 하고 계신 것", 14, bold=True, color=BRAND_BLUE)

    good_items = [
        "✓  리뷰 900건+ — 부산 내 신뢰도 확보 완료",
        "✓  블로그 체험단 15팀 — 기본 SEO 자산 축적",
        "✓  돼지우동 — 명확한 시그니처 메뉴 보유",
    ]
    tb = add_textbox(slide, 1.1, 2.6, 5, 2.0)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(good_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"
        p.space_after = Pt(8)

    # 우측 - 놓치고 있는 것
    add_rounded_rect(slide, 7.0, 2.0, 5.5, 2.8, WARM_BG)
    tb = add_textbox(slide, 7.3, 2.1, 5, 0.4)
    set_text(tb, "놓치고 계신 것", 14, bold=True, color=BRAND_ORANGE)

    bad_items = [
        "✗  월 100~150만원 투입 vs 서울 인지도 제로",
        "✗  인스타 \"기본적 내용만\" — 알고리즘 노출 안됨",
        "✗  동동국밥 서울 마곡 진출 → 선점 기회 감소 중",
    ]
    tb = add_textbox(slide, 7.3, 2.6, 5, 2.0)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bad_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"
        p.space_after = Pt(8)

    # 하단 핵심 메시지 박스
    add_rounded_rect(slide, 0.8, 5.2, 11.7, 1.5, BRAND_NAVY)
    tb = add_textbox(slide, 1.2, 5.35, 11, 1.2)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "좋은 메뉴(돼지우동)가 있는데, 부산 밖에서는 아무도 모릅니다."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "동동국밥이 서울에 먼저 깃발 꽂으면 \"국밥 프랜차이즈 = 동동\" 이미지가 굳어집니다."
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_ORANGE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 3: 핵심 문제 — 왜 지금인가
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 3, 0.35)
    set_text(tb, "02  핵심 문제", 12, bold=True, color=BRAND_BLUE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "지금이 아니면, 6개월 후엔 3배 비용이 듭니다", 30, bold=True, color=BLACK)

    # 3개 타이밍 카드
    timings = [
        ("① 동동국밥\n서울 진출", "서울에서 \"국밥 프랜차이즈\"\n하면 떠오르는 브랜드가\n아직 없습니다.\n\n동동국밥이 자리 잡기 전에\n온라인을 선점해야 합니다.", BRAND_BLUE),
        ("② 홈페이지\n3월 오픈", "홈페이지가 생기면 드디어\n온라인 마케팅의 \"도착지\"가\n완성됩니다.\n\n광고 → 홈페이지\n→ 가맹 문의 퍼널 완성.", BRAND_BLUE),
        ("③ 돼지우동\n바이럴 적기", "숏폼 콘텐츠 전성시대.\n\"국밥집 우동\"이라는 의외성은\n릴스/숏츠에 최적화된\n소재입니다.\n\n지금이 딱 맞는 타이밍.", BRAND_ORANGE),
    ]
    for idx, (title, desc, color) in enumerate(timings):
        cx = 0.8 + idx * 4.1
        card = add_rounded_rect(slide, cx, 2.0, 3.7, 4.5, WHITE, border_color=color, border_width=2)
        # 상단 색상 바
        add_rect(slide, cx, 2.0, 3.7, 0.08, color)
        # 제목
        tb = add_textbox(slide, cx + 0.3, 2.3, 3.1, 0.8)
        set_text(tb, title, 18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        # 설명
        tb = add_textbox(slide, cx + 0.3, 3.3, 3.1, 3.0)
        set_text(tb, desc, 11, color=DARK_GRAY)

    # 하단 핵심
    add_rounded_rect(slide, 0.8, 6.7, 11.7, 0.6, BRAND_NAVY)
    tb = add_textbox(slide, 1.2, 6.75, 11, 0.5)
    set_text(tb, "지금 월 500만원으로 할 수 있는 일을, 6개월 후에는 1,500만원을 써도 못합니다.", 14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # 슬라이드 4: A안 — 바이럴 집중형
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 3, 0.35)
    set_text(tb, "03  A안", 12, bold=True, color=BRAND_BLUE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "바이럴 집중형 — 빠른 임팩트", 30, bold=True, color=BLACK)

    tb = add_textbox(slide, 0.8, 1.6, 11.7, 0.4)
    set_text(tb, "\"국밥집인데 우동이 제일 잘 나가요\" — 이 의외성이 바이럴의 핵심", 14, color=BRAND_ORANGE)

    months_a = [
        ("1개월차", "폭탄 투하", ["숏폼 영상 4편 제작", "인플루언서 3명 (5~20만)", "인스타 광고 (서울 2030)", "콘텐츠 전략 수립"]),
        ("2개월차", "확산", ["유튜버 2팀 (10~50만)", "숏폼 4편 (고객 반응 등)", "인스타/유튜브 광고 확대", "인스타 콘텐츠 운영 주3회"]),
        ("3개월차", "수확", ["가맹 LP 광고 + 리타겟팅", "성과 콘텐츠 제작", "서울 진출 상권 분석", "전략 컨설팅"]),
    ]
    for idx, (month, subtitle, items) in enumerate(months_a):
        cx = 0.8 + idx * 4.1
        add_rounded_rect(slide, cx, 2.2, 3.7, 3.8, WHITE, border_color=BRAND_BLUE, border_width=1.5)
        # 상단 바
        add_rect(slide, cx, 2.2, 3.7, 0.06, BRAND_BLUE)
        # 월
        tb = add_textbox(slide, cx + 0.2, 2.4, 3.3, 0.35)
        set_text(tb, f"{month}: {subtitle}", 14, bold=True, color=BRAND_BLUE)
        # 항목들
        tb = add_textbox(slide, cx + 0.2, 2.9, 3.3, 2.4)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "맑은 고딕"
            p.space_after = Pt(6)
        # 월 비용
        tb = add_textbox(slide, cx + 0.2, 5.3, 3.3, 0.4)
        set_text(tb, "월 500만원", 16, bold=True, color=BRAND_BLUE, alignment=PP_ALIGN.RIGHT)

    # KPI 바
    add_rounded_rect(slide, 0.8, 6.3, 11.7, 0.9, ACCENT_BLUE_BG)
    tb = add_textbox(slide, 1.2, 6.35, 11, 0.8)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "3개월 핵심 KPI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.font.name = "맑은 고딕"
    kpis = "릴스 총 조회 50만+   |   팔로워 5,000+   |   가맹 문의 15건+   |   예상 ROI 133~200%"
    p2 = tf.add_paragraph()
    p2.text = kpis
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_BLUE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 5: B안 — 브랜드+바이럴 ★추천
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 3, 0.35)
    set_text(tb, "04  B안 ★추천", 12, bold=True, color=BRAND_ORANGE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "브랜드 빌딩 + 바이럴 병행형", 30, bold=True, color=BLACK)

    tb = add_textbox(slide, 0.8, 1.6, 11.7, 0.4)
    set_text(tb, "바이럴로 \"관심\" → 브랜드로 \"신뢰\" → 가맹으로 \"전환\"", 14, color=BRAND_ORANGE)

    months_b = [
        ("1개월차", "기초 공사", ["브랜드 전략 + 톤앤매너", "인스타 리브랜딩 (피드 9장)", "숏폼 3편 + 인플루언서 2명", "홈페이지 브랜드 감수"]),
        ("2개월차", "콘텐츠 체계화", ["콘텐츠 캘린더 (주 3회)", "브랜드 스토리 시리즈", "유튜버 1팀 + 광고 확대", "숏폼 3편 제작"]),
        ("3개월차", "가맹 전환", ["가맹 안내 콘텐츠 + LP 광고", "네이버 SEO 공략 5편", "리타겟팅 광고", "서울 진출 컨설팅"]),
    ]
    for idx, (month, subtitle, items) in enumerate(months_b):
        cx = 0.8 + idx * 4.1
        add_rounded_rect(slide, cx, 2.2, 3.7, 3.8, WHITE, border_color=BRAND_ORANGE, border_width=1.5)
        add_rect(slide, cx, 2.2, 3.7, 0.06, BRAND_ORANGE)
        tb = add_textbox(slide, cx + 0.2, 2.4, 3.3, 0.35)
        set_text(tb, f"{month}: {subtitle}", 14, bold=True, color=BRAND_ORANGE)
        tb = add_textbox(slide, cx + 0.2, 2.9, 3.3, 2.4)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "맑은 고딕"
            p.space_after = Pt(6)
        tb = add_textbox(slide, cx + 0.2, 5.3, 3.3, 0.4)
        set_text(tb, "월 500만원", 16, bold=True, color=BRAND_ORANGE, alignment=PP_ALIGN.RIGHT)

    # 핵심 차이 바
    add_rounded_rect(slide, 0.8, 6.3, 11.7, 0.9, WARM_BG)
    tb = add_textbox(slide, 1.2, 6.35, 11, 0.8)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "핵심 차이"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE
    p.font.name = "맑은 고딕"
    p2 = tf.add_paragraph()
    p2.text = "계약 종료 후에도 브랜드 자산(가이드라인, 콘텐츠 50개+, 팔로워)이 남습니다"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_ORANGE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 6: A vs B 비교
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 5, 0.35)
    set_text(tb, "05  A안 vs B안 비교", 12, bold=True, color=BRAND_BLUE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "같은 예산, 다른 전략. 어떤 것이 맞을까요?", 30, bold=True, color=BLACK)

    # 비교 테이블
    rows_data = [
        ("비교 항목", "A안: 바이럴 집중", "B안: 브랜드+바이럴 ★"),
        ("월 비용", "500만원", "500만원"),
        ("첫 달 임팩트", "★★★★★", "★★★"),
        ("3개월 후 자산", "★★", "★★★★★"),
        ("가맹 문의 속도", "빠름", "중간"),
        ("가맹 전환율", "중간", "높음"),
        ("계약 종료 후", "효과 급감", "자산 유지"),
        ("리스크", "바이럴 실패 시 성과 없음", "바이럴 부진해도 자산 남음"),
    ]

    col_widths = [Inches(3.0), Inches(4.0), Inches(4.0)]
    table_left = Inches(1.2)
    table_top = Inches(1.9)
    row_height = Inches(0.48)

    for ri, (c1, c2, c3) in enumerate(rows_data):
        y = 1.9 + ri * 0.52
        is_header = ri == 0
        bg = BRAND_NAVY if is_header else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        txt_color = WHITE if is_header else DARK_GRAY

        # 3 columns
        for ci, (text, w_inch) in enumerate(zip([c1, c2, c3], [3.0, 4.0, 4.0])):
            x = 1.2 + sum([3.0, 4.0, 4.0][:ci])
            cell = add_rect(slide, x, y, w_inch, 0.48, bg)
            tb = add_textbox(slide, x + 0.15, y + 0.05, w_inch - 0.3, 0.38)
            align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            set_text(tb, text, 11, bold=is_header, color=txt_color, alignment=align)

    # 추천
    tb = add_textbox(slide, 0.8, 6.4, 11.7, 0.7)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "→  B안 추천: 대표님이 \"500만원 어디에 썼어?\"라고 물으실 때"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE
    p.font.name = "맑은 고딕"
    p2 = tf.add_paragraph()
    p2.text = "    숫자(바이럴 성과) + 자산(브랜드 가이드라인) 둘 다 보여드릴 수 있습니다."
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_ORANGE
    p2.font.name = "맑은 고딕"

    # ═══════════════════════════════════════════
    # 슬라이드 7: 예상 ROI
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 3, 0.35)
    set_text(tb, "06  예상 ROI", 12, bold=True, color=BRAND_BLUE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "가맹 1건이면 마케팅비 2개월분을 회수합니다", 30, bold=True, color=BLACK)

    roi_cards = [
        ("보수적", "가맹 문의 10~15건\n실제 계약 1~2건\n가맹비 수익\n1,000~2,000만원", "ROI\n67~133%", MID_GRAY),
        ("기대치", "가맹 문의 15~20건\n실제 계약 3건\n가맹비 수익\n3,000만원", "ROI\n200%", BRAND_BLUE),
        ("낙관적", "가맹 문의 20건+\n실제 계약 4~5건\n가맹비 수익\n4,000~5,000만원", "ROI\n333%+", BRAND_ORANGE),
    ]

    for idx, (title, desc, roi, color) in enumerate(roi_cards):
        cx = 0.8 + idx * 4.1
        add_rounded_rect(slide, cx, 1.8, 3.7, 4.0, WHITE, border_color=color, border_width=2)
        add_rect(slide, cx, 1.8, 3.7, 0.08, color)
        # 타이틀
        tb = add_textbox(slide, cx + 0.3, 2.0, 3.1, 0.4)
        set_text(tb, title, 16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        # 상세
        tb = add_textbox(slide, cx + 0.3, 2.5, 3.1, 1.8)
        set_text(tb, desc, 11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
        # ROI 숫자
        tb = add_textbox(slide, cx + 0.3, 4.3, 3.1, 1.0)
        set_text(tb, roi, 24, bold=True, color=color, alignment=PP_ALIGN.CENTER)

    # 하단 핵심
    add_rounded_rect(slide, 0.8, 6.1, 11.7, 1.1, BRAND_NAVY)
    tb = add_textbox(slide, 1.2, 6.2, 11, 0.9)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "3개월 투자 1,500만원  →  가맹 1건(1,000만원)이면 67% 회수  →  2건이면 100%+ 회수"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "마케팅이 아니라 투자입니다."
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_ORANGE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 8: 브랜드라이즈 소개
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    tb = add_textbox(slide, 0.8, 0.4, 4, 0.35)
    set_text(tb, "07  브랜드라이즈", 12, bold=True, color=BRAND_BLUE)

    tb = add_textbox(slide, 0.8, 0.9, 11.7, 0.7)
    set_text(tb, "전략부터 실행까지, 한 팀이 책임집니다", 30, bold=True, color=BLACK)

    services = [
        ("전략 설계", "바이럴 컨셉 설계\n타겟 설정, KPI 관리"),
        ("콘텐츠 제작", "숏폼 영상 기획·디렉팅\n인스타 피드·스토리 기획"),
        ("인플루언서", "섭외·협상·가이드라인\n성과 관리·ROI 분석"),
        ("광고 운영", "메타·유튜브 광고\n세팅·최적화·리포팅"),
        ("브랜드 관리", "톤앤매너·비주얼 가이드\n홈페이지 연계·감수"),
        ("월간 리포팅", "성과 보고·전략 조정\n대표님 보고용 요약"),
    ]
    for idx, (svc_title, svc_desc) in enumerate(services):
        row = idx // 3
        col = idx % 3
        cx = 0.8 + col * 4.1
        cy = 1.9 + row * 2.2
        add_rounded_rect(slide, cx, cy, 3.7, 1.8, CARD_BG)
        # 아이콘 대신 제목에 색상
        tb = add_textbox(slide, cx + 0.3, cy + 0.15, 3.1, 0.4)
        set_text(tb, svc_title, 15, bold=True, color=BRAND_BLUE)
        tb = add_textbox(slide, cx + 0.3, cy + 0.65, 3.1, 1.0)
        set_text(tb, svc_desc, 11, color=DARK_GRAY)

    # 하단
    add_rounded_rect(slide, 0.8, 6.5, 11.7, 0.7, ACCENT_BLUE_BG)
    tb = add_textbox(slide, 1.2, 6.55, 11, 0.6)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "변미혜 팀장님은 현장 운영과 인플루언서 방문 응대에만 집중하세요."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BRAND_NAVY
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "온라인 전략과 실행은 저희가 맡겠습니다."
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_BLUE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 9: Next Step
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, BRAND_NAVY)

    # 상단 라인
    add_rect(slide, 3, 1.5, 7.333, 0.05, BRAND_ORANGE)

    tb = add_textbox(slide, 2, 1.8, 9.333, 0.8)
    set_text(tb, "Next Step", 42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    steps = [
        "①  A안 / B안 선택",
        "②  대표님 최종 승인",
        "③  계약 체결 → 1주일 내 킥오프",
        "④  1개월차 실행 → 즉시 콘텐츠 제작 시작",
    ]
    tb = add_textbox(slide, 3, 3.0, 7.333, 2.5)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"
        p.space_after = Pt(14)

    # 하단 라인
    add_rect(slide, 3, 5.7, 7.333, 0.05, BRAND_ORANGE)

    tb = add_textbox(slide, 2, 5.9, 9.333, 0.8)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "브랜드라이즈  |  brandrise.co.kr"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════
    # 슬라이드 10: 1-page 대표 보고용 요약
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    # 상단 바
    add_rect(slide, 0, 0, 13.333, 0.6, BRAND_BLUE)
    tb = add_textbox(slide, 0.5, 0.1, 12, 0.4)
    set_text(tb, "미례국밥 마케팅 제안 — 핵심 요약 (대표님 보고용)", 14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # A안 카드 (좌)
    add_rounded_rect(slide, 0.8, 1.0, 5.6, 4.8, WHITE, border_color=BRAND_BLUE, border_width=2)
    tb = add_textbox(slide, 1.1, 1.2, 5, 0.45)
    set_text(tb, "A안: 바이럴 집중형", 18, bold=True, color=BRAND_BLUE, alignment=PP_ALIGN.CENTER)

    a_body = "월 500만원 × 3개월 = 1,500만원"
    tb = add_textbox(slide, 1.3, 1.8, 4.6, 0.4)
    set_text(tb, a_body, 12, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    a_items = [
        "• 돼지우동 숏폼 바이럴 (릴스/숏츠)",
        "• 인플루언서 5명 + 유튜버 2팀",
        "• SNS 광고 집중 (서울 2030 타겟)",
        "",
        "장점: 빠른 성과, 즉각적 임팩트",
        "단점: 계약 종료 시 효과 급감",
    ]
    tb = add_textbox(slide, 1.3, 2.3, 4.6, 2.5)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(a_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"
        p.space_after = Pt(4)

    # A안 ROI
    add_rounded_rect(slide, 1.3, 4.6, 4.6, 0.8, ACCENT_BLUE_BG)
    tb = add_textbox(slide, 1.5, 4.65, 4.2, 0.7)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "예상 ROI: 133~200%"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "(보수적 기준)"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # B안 카드 (우)
    add_rounded_rect(slide, 6.9, 1.0, 5.6, 4.8, WHITE, border_color=BRAND_ORANGE, border_width=2)
    tb = add_textbox(slide, 7.2, 1.2, 5, 0.45)
    set_text(tb, "B안: 브랜드+바이럴 ★추천", 18, bold=True, color=BRAND_ORANGE, alignment=PP_ALIGN.CENTER)

    b_body = "월 500만원 × 3개월 = 1,500만원"
    tb = add_textbox(slide, 7.4, 1.8, 4.6, 0.4)
    set_text(tb, b_body, 12, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    b_items = [
        "• 브랜드 전략 + 톤앤매너 정립",
        "• 돼지우동 바이럴 (숏폼 + 광고)",
        "• 콘텐츠 시스템 구축 (주 3회)",
        "",
        "장점: 자산 축적, 가맹 전환율 높음",
        "단점: 첫 달 임팩트 상대적 약함",
    ]
    tb = add_textbox(slide, 7.4, 2.3, 4.6, 2.5)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(b_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"
        p.space_after = Pt(4)

    # B안 ROI
    add_rounded_rect(slide, 7.4, 4.6, 4.6, 0.8, WARM_BG)
    tb = add_textbox(slide, 7.6, 4.65, 4.2, 0.7)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "예상 ROI: 200%+"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "(기대치 기준)"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    # 하단 핵심 메시지 (공통)
    add_rounded_rect(slide, 0.8, 6.1, 11.7, 1.1, BRAND_NAVY)
    tb = add_textbox(slide, 1.2, 6.2, 11, 0.9)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "가맹 1건 (가맹비 1,000만원) = 마케팅비 2개월분 즉시 회수"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "동동국밥 서울 진출 전에 온라인을 선점해야 합니다"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = BRAND_ORANGE
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER

    return prs


def upload_to_drive(filepath, title):
    """Google Drive에 PPTX 업로드 후 Google Slides로 변환"""
    try:
        from google.oauth2.credentials import Credentials as OAuthCreds
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload

        token_path = os.path.join(os.path.dirname(__file__), "..", ".secrets", "token.json")
        if not os.path.exists(token_path):
            token_path = os.path.expanduser("~/AI/.secrets/token.json")

        with open(token_path) as f:
            token_data = json.load(f)
        creds = OAuthCreds(
            token=token_data["token"],
            refresh_token=token_data["refresh_token"],
            token_uri=token_data["token_uri"],
            client_id=token_data["client_id"],
            client_secret=token_data["client_secret"],
            scopes=token_data["scopes"],
        )
        drive = build("drive", "v3", credentials=creds)

        file_metadata = {
            "name": title,
            "mimeType": "application/vnd.google-apps.presentation",
        }
        media = MediaFileUpload(
            filepath,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            resumable=True,
        )
        file = drive.files().create(
            body=file_metadata, media_body=media, fields="id"
        ).execute()

        file_id = file.get("id")
        url = f"https://docs.google.com/presentation/d/{file_id}/edit"
        print(f"📎 Google Slides URL: {url}")
        return url
    except Exception as e:
        print(f"⚠️  Google Drive 업로드 실패: {e}")
        print(f"📁 로컬 파일로 사용하세요: {filepath}")
        return None


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    print("슬라이드 생성 중...")
    prs = build_presentation()
    prs.save(OUTPUT_FILE)
    print(f"✅ PPTX 저장 완료: {OUTPUT_FILE}")

    print("\nGoogle Drive에 업로드 중...")
    url = upload_to_drive(OUTPUT_FILE, "미례국밥 마케팅 제안서 — 브랜드라이즈")

    if url:
        print(f"\n🎉 완료! Google Slides에서 확인하세요:")
        print(f"   {url}")
    else:
        print(f"\n📁 로컬 PPTX 파일: {OUTPUT_FILE}")
        print("   Google Drive에 수동 업로드하시면 Google Slides로 변환됩니다.")


if __name__ == "__main__":
    main()
