"""Export brandrise-company-profile-v2.html to PPTX."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_CREAM = RGBColor(0xFA, 0xF6, 0xF1)
BG_WARM = RGBColor(0xF5, 0xEF, 0xE8)
ORANGE = RGBColor(0xE8, 0x65, 0x2B)
ORANGE_LIGHT = RGBColor(0xFD, 0xF0, 0xEA)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SUB = RGBColor(0x55, 0x55, 0x55)
TEXT_MUTED = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE8, 0xE0, 0xD6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=BG_CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align

    # Handle multiline
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_card(slide, left, top, width, height, fill_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_footer(slide, page_num):
    add_text(slide, 0.6, 6.9, 3, 0.4, '브랜드라이즈 서비스 소개서', size=8, color=TEXT_MUTED)
    add_text(slide, 11.5, 6.9, 1.5, 0.4, str(page_num), size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

def add_nav(slide, current, total=10, label='브랜드라이즈 소개'):
    # Nav background
    nav_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.65))
    nav_bg.fill.solid()
    nav_bg.fill.fore_color.rgb = BG_WARM
    nav_bg.line.fill.background()

    add_text(slide, 3.0, 0.15, 2.5, 0.35, label, size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

    # Number chips
    start_x = 5.8
    for i in range(1, total + 1):
        x = start_x + (i - 1) * 0.45
        if i == current:
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.13), Inches(0.38), Inches(0.38))
            chip.fill.solid()
            chip.fill.fore_color.rgb = ORANGE
            chip.line.fill.background()
            tf = chip.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i)
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = WHITE
        else:
            add_text(slide, x, 0.15, 0.38, 0.35, str(i), size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# P1. COVER
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

add_text(slide, 0.8, 1.5, 3, 0.4, 'BRANDRISE', size=10, color=TEXT_MUTED)
add_text(slide, 0.8, 2.5, 9, 1.5, '마케팅 비용은 쓰고 있는데\n브랜드는 안 쌓이고 있다면.', size=40, bold=True, color=TEXT_DARK)
# Orange highlight on second line
add_text(slide, 0.8, 4.2, 9, 0.8, '직접 브랜드를 만들어 3개국에 진출시킨 팀이\n당신의 스타트업에 맞는 브랜딩 전략을 설계합니다.', size=14, color=TEXT_SUB)

# Tag
tag = add_card(slide, 0.8, 5.3, 2.8, 0.45)
tag.fill.solid()
tag.fill.fore_color.rgb = WHITE
add_text(slide, 0.85, 5.33, 2.7, 0.4, '   스타트업 브랜드 전략 파트너', size=10, color=TEXT_SUB)
# Orange dot
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(5.47), Inches(0.08), Inches(0.08))
dot.fill.solid()
dot.fill.fore_color.rgb = ORANGE
dot.line.fill.background()

add_text(slide, 0.6, 6.9, 4, 0.4, '브랜드라이즈 서비스 소개서', size=8, color=TEXT_MUTED)
add_text(slide, 9.0, 6.9, 4, 0.4, '주식회사 히즈 (HIZ Inc.)', size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════
# P2. PROBLEM
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 1)

add_text(slide, 0.8, 1.0, 8, 0.8, '왜 스타트업 마케팅은\n실패할까요?', size=30, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 2.2, 7, 0.9, '매달 200만원의 광고비를 쓰고 있지만,\n마케팅이 안 되고 있습니다.\n어디부터 손대야 하는지 모르겠다면,\n가장 먼저 해야 할 것은 광고가 아닙니다.', size=13, color=TEXT_SUB)

# 3 cards
cards_data = [
    ('73%', '전략 없이 집행', '초기 스타트업 마케팅 예산의\n대부분이 브랜드 자산으로\n남지 않고 소진됩니다.'),
    ('6개월', '평균 이탈 시점', '대행사를 바꾸거나 마케팅을\n중단하는 평균 시점.\n"효과를 모르겠다"가 주된 이유.'),
    ('2.3배', 'CAC 격차', '브랜딩이 된 스타트업과\n안 된 스타트업의\n고객획득비용 차이입니다.'),
]
for i, (num, title, desc) in enumerate(cards_data):
    x = 0.8 + i * 4.0
    add_card(slide, x, 3.5, 3.6, 2.5)
    add_text(slide, x + 0.3, 3.7, 3, 0.5, num, size=24, bold=True, color=ORANGE)
    add_text(slide, x + 0.3, 4.3, 3, 0.3, title, size=13, bold=True, color=TEXT_DARK)
    add_text(slide, x + 0.3, 4.7, 3, 1.0, desc, size=10, color=TEXT_SUB)

# Bottom bar
bottom = add_card(slide, 0.8, 6.2, 11.7, 0.55, ORANGE_LIGHT)
bottom.line.fill.background()
add_text(slide, 1.0, 6.25, 11.3, 0.45, '마케팅이 문제가 아닙니다.  브랜딩 없이 마케팅만 하는 것이 문제입니다.', size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_footer(slide, 2)


# ═══════════════════════════════════════════════
# P3. SOLUTION
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 2)

add_text(slide, 0.8, 1.0, 9, 0.8, '브랜드라이즈는\n이 문제를 이렇게 풀어드립니다.', size=30, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 2.2, 7, 0.4, '전략만 세워주고 떠나는 컨설팅이 아닙니다. 직접 실행하고, 숫자로 증명합니다.', size=13, color=TEXT_SUB)

sol_data = [
    ('01', '진단', '현재 브랜드 상태를 정밀 분석하고\n경쟁사 대비 포지셔닝을 파악합니다.'),
    ('02', '솔루션', '브랜드 상황에 맞는\n맞춤 마케팅 액션을 설계합니다.'),
    ('03', '분석(Report)', '매출/전환 데이터를 분석하고\n전략을 지속적으로 조정합니다.'),
]
for i, (num, title, desc) in enumerate(sol_data):
    x = 0.8 + i * 4.0
    add_card(slide, x, 3.2, 3.6, 3.0)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1.3), Inches(3.5), Inches(1.0), Inches(0.55), )
    circ.fill.solid()
    circ.fill.fore_color.rgb = ORANGE_LIGHT
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ORANGE

    add_text(slide, x + 0.3, 4.3, 3, 0.4, title, size=15, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 4.8, 3, 1.0, desc, size=10, color=TEXT_SUB, align=PP_ALIGN.CENTER)
add_footer(slide, 3)


# ═══════════════════════════════════════════════
# P4. KEY METRICS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 3)

add_text(slide, 0.8, 1.0, 10, 0.8, '브랜딩 & 마케팅 프로페셔널\n18년의 노하우를 스타트업에.', size=30, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 2.2, 7, 0.4, '대기업에서 검증된 전략 체계를 스타트업 규모에 맞게 최적화합니다.', size=13, color=TEXT_SUB)

metrics = [
    ('프로젝트 경험', '30+', '대기업·스타트업 브랜드 프로젝트'),
    ('평균 연속 수주', '3년', '클라이언트 재계약 기간'),
    ('글로벌 진출', '3개국', '자사 브랜드 미국·일본·한국 운영'),
]
for i, (label, value, desc) in enumerate(metrics):
    x = 0.8 + i * 4.0
    add_card(slide, x, 3.2, 3.6, 2.8)
    add_text(slide, x + 0.3, 3.5, 3, 0.3, label, size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 4.0, 3, 0.7, value, size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 4.9, 3, 0.3, desc, size=10, color=TEXT_SUB, align=PP_ALIGN.CENTER)
add_footer(slide, 4)


# ═══════════════════════════════════════════════
# P5. SERVICES
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 4, label='서비스 영역')

add_text(slide, 0.8, 0.85, 1.5, 0.25, 'SERVICES', size=8, bold=True, color=ORANGE)
add_text(slide, 0.8, 1.1, 8, 0.5, '브랜드라이즈가 만들면 다릅니다', size=28, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 1.7, 5, 0.3, '필요한 영역만 선택하여 진행할 수 있습니다.', size=12, color=TEXT_SUB)

svcs = [
    ('01', '브랜드 전략', 'Brand Strategy', '시장 분석부터 포지셔닝까지\n브랜드가 설 자리를 설계합니다.', '경쟁사 분석 & 시장 기회 발굴\n포지셔닝 & 핵심 메시지\nGTM 로드맵'),
    ('02', 'BI / CI 디자인', 'Brand Identity', '로고·컬러·톤앤매너·패키지\n브랜드만의 목소리를 만듭니다.', '비주얼 아이덴티티 시스템\n브랜드 가이드라인\n패키지 & 굿즈 디자인'),
    ('03', '채널 전략', 'Channel Strategy', 'SNS·유통·퍼포먼스 통합 설계\n팔리는 구조를 만듭니다.', '퍼포먼스 마케팅 (Meta, Google)\nSNS 콘텐츠 전략\n유통 채널 최적화'),
    ('04', '외부 CMO 팀', 'CMO as a Service', '월간 브랜드 코칭 & 성과 관리\n전략부터 실행까지 함께합니다.', '주간/월간 브랜드 코칭\n대행사 핸들링 & 감독\n통합 KPI 리포트'),
    ('05', '글로벌 진출', 'Global Expansion', '미국·일본 시장 진출 전략\n831개 해외 바이어 네트워크.', '해외 마켓 진입 전략\nCosmoprof 전시 기획\n현지 유통 네트워크'),
    ('06', 'AI 브랜드 가이드', 'AI Brand Guide', 'AI 활용 운영 로드맵 설계\n혼자서도 운영할 수 있는 시스템.', 'AI 마케팅 자동화 설계\n콘텐츠 생산 시스템\n운영 매뉴얼 & 교육'),
]
for i, (num, title, sub, desc, details) in enumerate(svcs):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 2.2 + row * 2.7
    add_card(slide, x, y, 3.6, 2.4)
    add_text(slide, x + 0.25, y + 0.15, 3, 0.2, f'Service {num}', size=8, bold=True, color=ORANGE)
    add_text(slide, x + 0.25, y + 0.4, 3, 0.3, title, size=13, bold=True, color=TEXT_DARK)
    add_text(slide, x + 0.25, y + 0.7, 3, 0.2, sub, size=9, color=ORANGE)
    add_text(slide, x + 0.25, y + 1.0, 3.1, 0.5, desc, size=9, color=TEXT_SUB)
    # Details with dots
    detail_lines = '\n'.join(['· ' + d for d in details.split('\n')])
    add_text(slide, x + 0.25, y + 1.55, 3.1, 0.8, detail_lines, size=8, color=TEXT_SUB)
add_footer(slide, 5)


# ═══════════════════════════════════════════════
# P6. PROCESS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 5, label='진행 프로세스')

add_text(slide, 0.8, 0.85, 1.5, 0.25, 'PROCESS', size=8, bold=True, color=ORANGE)
add_text(slide, 0.8, 1.1, 8, 0.5, '4주 안에 브랜드의 방향이 잡힙니다.', size=28, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 1.7, 7, 0.4, '긴 보고서 대신, 빠른 실행. 4주 집중 프로그램으로 브랜드의 뾰족함을 만듭니다.', size=12, color=TEXT_SUB)

steps = [
    ('Week 1', '진단', '현재 위치 파악\n경쟁사 분석\n핵심 과제 도출'),
    ('Week 2', '발견', '뾰족함 도출\n타겟 고객 정의\n포지셔닝 확정'),
    ('Week 3', '설계', '전략 문서화\n실행 로드맵\n콘텐츠 기획안'),
    ('Week 4', '런칭', '첫 실행 착수\n성과 측정 세팅\n주간 케어 전환'),
]
for i, (week, title, desc) in enumerate(steps):
    x = 0.8 + i * 3.1
    add_card(slide, x, 2.5, 2.8, 3.0)
    add_text(slide, x + 0.2, 2.7, 2.4, 0.25, week, size=8, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 3.05, 2.4, 0.35, title, size=16, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 3.5, 2.4, 1.2, desc, size=10, color=TEXT_SUB, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 3:
        add_text(slide, x + 2.7, 3.6, 0.5, 0.4, '→', size=16, color=ORANGE, align=PP_ALIGN.CENTER)

# Note
note = add_card(slide, 0.8, 5.8, 11.7, 0.55, ORANGE_LIGHT)
note.line.fill.background()
add_text(slide, 1.0, 5.85, 11.3, 0.4, '4주 이후 → 주간 케어 모델로 전환하여 지속적으로 성장을 함께합니다.', size=11, color=TEXT_SUB, align=PP_ALIGN.CENTER)
add_footer(slide, 6)


# ═══════════════════════════════════════════════
# P7. WHY US
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 6, label='Why Brandrise')

add_text(slide, 0.8, 0.85, 2, 0.25, 'WHY BRANDRISE', size=8, bold=True, color=ORANGE)
add_text(slide, 0.8, 1.1, 10, 0.5, '직접 브랜드를 만들어본 사람이,\n대기업 급 전략을 씁니다.', size=28, bold=True, color=TEXT_DARK)
add_text(slide, 0.8, 2.0, 7, 0.4, '남의 브랜드만 키운 게 아닙니다. 내 돈으로, 내 브랜드를, 바닥부터 만들어봤습니다.', size=12, color=TEXT_SUB)

# Left card - Startup
add_card(slide, 0.8, 2.8, 5.7, 3.0)
# Orange top line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(5.7), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ORANGE
line.line.fill.background()
add_text(slide, 1.1, 2.95, 5, 0.2, '직접 브랜드를 만들어 성공시킨 경험', size=8, bold=True, color=ORANGE)
add_text(slide, 1.1, 3.2, 5, 0.4, 'ReSaltZ — 0에서 시작해\n미국/일본 진출까지', size=15, bold=True, color=TEXT_DARK)
startup_items = '· 프리미엄 웰니스 브랜드 기획부터 A-to-Z\n· 미국 최상위 마켓 진입 + 일본 유통 확보\n· 뉴욕, 도쿄, 서울 — 3개국 팝업 직접 운영\n· 출시 18개월 만에 연 매출 8억 돌파\n· 미국 법인 설립 (Seoul Wellness, INC.)'
add_text(slide, 1.1, 3.8, 5, 1.8, startup_items, size=10, color=TEXT_SUB)

# Right card - Enterprise
add_card(slide, 6.8, 2.8, 5.7, 3.0)
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.8), Inches(5.7), Inches(0.04))
line2.fill.solid()
line2.fill.fore_color.rgb = ORANGE
line2.line.fill.background()
add_text(slide, 7.1, 2.95, 5, 0.2, '18년 기업 브랜드 노하우', size=8, bold=True, color=ORANGE)
add_text(slide, 7.1, 3.2, 5, 0.4, '글로벌 브랜드 IMC 전략\n— 전략에서 캠페인까지 풀스택', size=15, bold=True, color=TEXT_DARK)
corp_items = '· 아모레퍼시픽, CJ제일제당, 네슬레\n· 하이네켄, 스타벅스엔홈, 현대차정몽구재단\n· ETUDE, MISSHA, 현대백화점그룹\n· Cosmoprof NA/Bologna 직접 참가 및 수주\n· 대기업 팝업 30건+ 기획/운영 실행'
add_text(slide, 7.1, 3.8, 5, 1.8, corp_items, size=10, color=TEXT_SUB)

# Bottom
bottom = add_card(slide, 0.8, 6.0, 11.7, 0.6, ORANGE_LIGHT)
bottom.line.fill.background()
add_text(slide, 1.0, 6.05, 11.3, 0.5, '에이전시는 남의 브랜드를 키웁니다.  브랜드라이즈는 직접 키워본 사람이 만듭니다.', size=13, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_footer(slide, 7)


# ═══════════════════════════════════════════════
# P8. CLIENTS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 7, label='고객사')

add_text(slide, 0.8, 0.85, 1.5, 0.25, 'CLIENTS', size=8, bold=True, color=ORANGE)
add_text(slide, 0.8, 1.1, 10, 0.5, '스타트업부터 대기업까지,\n함께 성장해온 파트너들.', size=28, bold=True, color=TEXT_DARK)

clients = [
    ('ReSaltZ', '프리미엄 웰니스'), ('VEGGIET', '식물성 푸드테크'), ('WEETAMIN', '건강기능식품'), ('CheckVO', '체크보'),
    ('하이네켄', 'Heineken'), ('CJ제일제당', 'CJ CheilJedang'), ('네슬레', 'Nestle'), ('아모레퍼시픽', 'Amorepacific'),
    ('스타벅스앳홈', 'Starbucks at Home'), ('현대백화점그룹', 'Hyundai Dept.'), ('ETUDE', '에뛰드'), ('MISSHA', '미샤'),
    ('요기요', 'Yogiyo'), ('NC', 'NCSoft'), ('정몽구재단', 'CMK Foundation'), ('현대글로비스', 'Hyundai Glovis'),
]
for i, (name, sub) in enumerate(clients):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.05
    y = 2.3 + row * 1.1
    add_card(slide, x, y, 2.8, 0.9)
    add_text(slide, x + 0.15, y + 0.1, 2.5, 0.35, name, size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, y + 0.5, 2.5, 0.25, sub, size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 5.5, 6.8, 2.5, 0.3, '외 다수 프로젝트', size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 8)


# ═══════════════════════════════════════════════
# P9. TEAM
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav(slide, 8, label='팀 소개')

add_text(slide, 0.8, 0.85, 1.5, 0.25, 'TEAM', size=8, bold=True, color=ORANGE)
add_text(slide, 0.8, 1.1, 8, 0.5, '누가 당신의 브랜드를\n함께 만드는가.', size=28, bold=True, color=TEXT_DARK)

# Team card
add_card(slide, 0.8, 2.3, 11.7, 4.3)

# Left - info
add_text(slide, 1.2, 2.5, 4, 0.35, '권윤정', size=20, bold=True, color=TEXT_DARK)
add_text(slide, 1.2, 2.95, 4, 0.25, '대표 · Brand Strategist', size=10, color=ORANGE)
add_text(slide, 1.2, 3.4, 4.5, 0.8, '18년간 글로벌 브랜드의 전략을 설계했습니다.\n그리고 직접 브랜드를 만들어 3개국에 진출시켰습니다.\n이제 그 경험을 당신의 브랜드에 씁니다.', size=10, color=TEXT_SUB)

career = '· 글로벌 브랜드 IMC 전략 18년+\n· 아모레퍼시픽, CJ, 네슬레, 하이네켄 외\n· ReSaltZ 브랜드 창업 → 미국/일본 진출\n· 대기업/자사 팝업 30건+ 직접 기획·운영\n· Cosmoprof NA/Bologna 참가 및 수주\n· 미국 법인 설립 (Seoul Wellness, INC.)\n· 서울 창창사 11기 · A-벤처스'
add_text(slide, 1.2, 4.3, 4.5, 2.0, career, size=9, color=TEXT_SUB)

# Right - quote
# Orange left border line
qline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.7), Inches(0.04), Inches(1.2))
qline.fill.solid()
qline.fill.fore_color.rgb = ORANGE
qline.line.fill.background()
add_text(slide, 6.8, 2.7, 5, 1.2, '"대행사에서 남의 브랜드를 키울 때는 몰랐습니다.\n내 돈으로 내 브랜드를 만들어보니,\n진짜 중요한 건 완전히 다르더라고요."', size=13, bold=True, color=TEXT_DARK)

add_text(slide, 6.8, 4.2, 5.3, 2.0, '대기업에서 배운 체계적인 전략.\n자사 브랜드에서 깨달은 현장의 감각.\n이 두 가지가 합쳐질 때, 스타트업의 브랜드는 비로소 뾰족해집니다.\n\n브랜드라이즈는 "예쁜 결과물"을 만드는 곳이 아닙니다.\n대표님이 밤새 고민하는 그 문제를, 함께 풀어가는 팀입니다.', size=10, color=TEXT_SUB)
add_footer(slide, 9)


# ═══════════════════════════════════════════════
# P10. CTA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WARM)
add_nav(slide, 9, label='Contact')

add_text(slide, 0.8, 0.85, 1.5, 0.25, 'CONTACT', size=8, bold=True, color=ORANGE)
add_text(slide, 1.5, 2.0, 10, 1.0, '아직도 브랜딩 플랜 없이\n마케팅 예산 쓰고 계신가요?', size=34, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.3, 8, 0.6, '무료 브랜드 진단으로 시작하세요.\n30분이면 지금 당장 개선할 수 있는 3가지를 찾아드립니다.', size=13, color=TEXT_SUB, align=PP_ALIGN.CENTER)

# CTA Button
btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(4.2), Inches(3.7), Inches(0.6))
btn.fill.solid()
btn.fill.fore_color.rgb = ORANGE
btn.line.fill.background()
tf = btn.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '무료 브랜드 진단 신청하기'
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

add_text(slide, 4.5, 4.9, 4.5, 0.3, '상담 후 진행 여부는 자유롭게 결정하세요.', size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Contact info
contacts = [('Instagram', '@brandrise_kr'), ('Email', 'hello@brandrise.kr'), ('Location', '서울 마포구 양화로')]
for i, (label, val) in enumerate(contacts):
    x = 3.2 + i * 2.5
    add_text(slide, x, 5.5, 2.2, 0.2, label, size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x, 5.75, 2.2, 0.3, val, size=10, color=TEXT_SUB, align=PP_ALIGN.CENTER)

add_text(slide, 2.5, 6.5, 8.5, 0.4, '주식회사 히즈 (HIZ Inc.) · 설립 2018 · 서울시 마포구 양화로 73-1, 5층\nBrand rise isn\'t luck. It\'s engineered.', size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 10)


# ── Save ──
output = '/Users/wooseongmin/AI/.claude/worktrees/new-business/content/brandrise-company-profile-v2.pptx'
prs.save(output)
print(f'Saved: {output}')
